import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    C_DARK = RGBColor(15, 23, 42)       # Slate 900
    C_WHITE = RGBColor(255, 255, 255)
    C_TEAL = RGBColor(13, 148, 136)     # Primary Teal
    C_TEAL_DARK = RGBColor(4, 47, 46)
    C_INDIGO = RGBColor(79, 70, 229)    # Predictive Accent
    C_CARD_BG = RGBColor(248, 250, 252) # Slate 50
    C_CARD_BORDER = RGBColor(226, 232, 240)
    C_TEXT_MUTED = RGBColor(100, 116, 139)
    C_TEXT_MAIN = RGBColor(30, 41, 59)
    C_CORAL = RGBColor(225, 29, 72)
    C_EMERALD = RGBColor(5, 150, 105)

    def add_header(slide, title_text, subtitle_text=""):
        # Header bar
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_TEAL
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Arial"
            p2.font.size = Pt(11)
            p2.font.color.rgb = C_TEXT_MUTED
            p2.space_before = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (High-Impact Hero)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Dark Hero Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_DARK
    bg1.line.color.rgb = C_DARK

    # Top Pill Badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.8), Inches(4.2), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_TEAL_DARK
    badge.line.color.rgb = C_TEAL
    badge_tf = badge.text_frame
    badge_p = badge_tf.paragraphs[0]
    badge_p.text = "VITISH 2026 • SIH INTERNAL HACKATHON"
    badge_p.alignment = PP_ALIGN.CENTER
    badge_p.font.name = "Arial"
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C_TEAL

    # Main Project Title
    title_box = slide1.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "MedFlow"
    p1.font.name = "Arial"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Real-Time Predictive Hospital Resource & Bed Turnover Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_TEAL
    p2.space_before = Pt(6)

    p3 = tf1.add_paragraph()
    p3.text = "Proactive 12–24h Inpatient Bed Turnover Engine • Specialty-Aware Referral • ABDM/FHIR • Rural USSD Fallback"
    p3.font.name = "Arial"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_before = Pt(6)

    # Info Cards Bottom Grid (3 Cards)
    card_data = [
        ("PROBLEM STATEMENT", "Hospital Bed & Resource Optimisation Platform\nCategory: Software | Theme: Healthcare / MedTech\nPS ID: [Insert PS ID]"),
        ("PRIMARY INNOVATION", "Predictive Bed Turnover Engine (ML)\nDischarge likelihood forecast 12–24h in advance\nReduces emergency hospital-hopping by 35–45%"),
        ("TEAM INFORMATION", "Team Name: [Insert Team Name]\nTeam Leader: [Insert Leader Name]\nInstitute: Vellore Institute of Technology (VIT)")
    ]

    for idx, (head, body) in enumerate(card_data):
        cx = Inches(0.9 + idx * 3.9)
        c = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(4.3), Inches(3.7), Inches(2.4))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(30, 41, 59)
        c.line.color.rgb = RGBColor(51, 65, 85)
        
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.2)
        ctf.margin_left = Inches(0.2)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = head
        cp1.font.name = "Arial"
        cp1.font.size = Pt(11)
        cp1.font.bold = True
        cp1.font.color.rgb = C_TEAL
        
        cp2 = ctf.add_paragraph()
        cp2.text = body
        cp2.font.name = "Arial"
        cp2.font.size = Pt(11)
        cp2.font.color.rgb = C_WHITE
        cp2.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 2: Idea Title & Proposed Solution
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "MedFlow — Proposed Solution & Value Proposition", "Three-Layer Systems Architecture solving healthcare information asymmetry across India")

    # Left Column: Proposed Solution & Differentiator (2 Cards)
    # Card 1: Core Solution
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = C_CARD_BG
    c1.line.color.rgb = C_CARD_BORDER
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_top = Inches(0.25)
    
    p = tf_c1.paragraphs[0]
    p.text = "PROPOSED SOLUTION"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_TEAL

    points_c1 = [
        ("Unified 3-Stakeholder Network: ", "Synchronizes Patients/Families, Hospital Staff, and District/State Administrators in real time."),
        ("Predictive Bed Turnover Engine: ", "Primary technical differentiator forecasting freed beds 12–24h ahead using ML Length-of-Stay modeling."),
        ("Smart Specialty-Aware Referral: ", "Multi-factor algorithm matching patients by specialty (e.g. Cardiac ICU vs Pediatric ICU), distance, and PMJAY."),
        ("Government Command Center: ", "Statewide GIS heatmap, automated shortage alerts (<10% ICU free), and inter-district reallocation."),
        ("Rural Inclusivity Fallback: ", "USSD (*999#) and SMS search for 2G feature-phone users in rural and low-connectivity regions.")
    ]
    for bold_txt, norm_txt in points_c1:
        p = tf_c1.add_paragraph()
        p.space_before = Pt(8)
        run_b = p.add_run()
        run_b.text = "• " + bold_txt
        run_b.font.bold = True
        run_b.font.size = Pt(11)
        run_b.font.color.rgb = C_TEXT_MAIN
        run_n = p.add_run()
        run_n.text = norm_txt
        run_n.font.size = Pt(11)
        run_n.font.color.rgb = C_TEXT_MAIN

    # Right Column: 3-Layer Impact Table + Uniqueness (Card 2)
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4))
    c2.fill.solid()
    c2.fill.fore_color.rgb = C_CARD_BG
    c2.line.color.rgb = C_CARD_BORDER
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_top = Inches(0.25)

    p = tf_c2.paragraphs[0]
    p.text = "THREE-LAYER PROBLEM & INNOVATION"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_INDIGO

    layers = [
        ("Patient Layer:", "Eliminates blind hospital hopping during Golden Hour with live bed, ICU, oxygen, and ambulance tracking."),
        ("Hospital Layer:", "Replaces manual registers with 1-click ward toggles, non-HMS shift handovers, and ML discharge readiness."),
        ("State Layer:", "Real-time aggregated telemetry prevents black-marketing, optimizes bulk oxygen, and coordinates mass casualties."),
        ("ABDM & FHIR R4:", "Native generation of interoperable HL7 FHIR R4 Document Bundles linked to 14-digit ABHA Health IDs."),
        ("Blockchain Audit Trail:", "Cryptographic SHA-256 hash chaining on all bed/oxygen changes for complete anti-tamper transparency.")
    ]
    for bold_txt, norm_txt in layers:
        p = tf_c2.add_paragraph()
        p.space_before = Pt(8)
        run_b = p.add_run()
        run_b.text = "✔ " + bold_txt + " "
        run_b.font.bold = True
        run_b.font.size = Pt(11)
        run_b.font.color.rgb = C_INDIGO
        run_n = p.add_run()
        run_n.text = norm_txt
        run_n.font.size = Pt(11)
        run_n.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------
    # SLIDE 3: Technical Approach & Architecture
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Technical Approach, Tech Stack & Workflow", "Modular microservices architecture, Scikit-Learn ML pipeline, and real-time WebSockets")

    # 4 Architecture Horizontal Cards
    stack_boxes = [
        ("FRONTEND LAYER", "React 18 + TypeScript (Vite)\nVanilla CSS Design System\nLeaflet GIS Route Mapping\nFlutter Mobile App / PWA", C_TEAL),
        ("BACKEND API & WS", "FastAPI (Python 3.13) / Node.js\nAsync REST Endpoints\nWebSockets (/ws/live) <100ms\nJWT Auth + Native Bcrypt", C_INDIGO),
        ("PREDICTIVE ML ENGINE", "Scikit-Learn Random Forest\nGradient Boosting LOS Regressor\n10 Clinical Feature Extraction\n12h & 24h Ward Capacity Rollup", C_EMERALD),
        ("DATABASE & STANDARDS", "PostgreSQL / SQLAlchemy ORM\nRedis Real-time Cache\nABDM APIs + HL7 FHIR R4\nDocker Compose & MeghRaj", C_TEAL_DARK)
    ]

    for idx, (title, body, col) in enumerate(stack_boxes):
        bx = Inches(0.8 + idx * 2.98)
        b = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(1.5), Inches(2.8), Inches(2.2))
        b.fill.solid()
        b.fill.fore_color.rgb = C_CARD_BG
        b.line.color.rgb = col
        b.line.width = Pt(1.5)
        btf = b.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_top = Inches(0.15)
        
        bp1 = btf.paragraphs[0]
        bp1.text = title
        bp1.font.name = "Arial"
        bp1.font.size = Pt(10.5)
        bp1.font.bold = True
        bp1.font.color.rgb = col
        
        bp2 = btf.add_paragraph()
        bp2.text = body
        bp2.font.name = "Arial"
        bp2.font.size = Pt(9.5)
        bp2.font.color.rgb = C_TEXT_MAIN
        bp2.space_before = Pt(4)

    # Workflow Bottom Card
    wf_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.733), Inches(3.0))
    wf_card.fill.solid()
    wf_card.fill.fore_color.rgb = RGBColor(15, 23, 42)
    wf_card.line.color.rgb = RGBColor(51, 65, 85)
    
    wf_tf = wf_card.text_frame
    wf_tf.word_wrap = True
    wf_tf.margin_left = wf_tf.margin_top = Inches(0.2)
    
    wfp1 = wf_tf.paragraphs[0]
    wfp1.text = "PRIMARY END-TO-END DEMO WORKFLOW"
    wfp1.font.name = "Arial"
    wfp1.font.size = Pt(12)
    wfp1.font.bold = True
    wfp1.font.color.rgb = RGBColor(45, 212, 191) # Light Teal

    flow_steps = [
        "1. Patient Searches Emergency Care → MedFlow evaluates real-time capacity and predicted 12h/24h bed gains.",
        "2. Specialty-Aware Referral Matcher → Computes Score: 35% Specialty + 25% Beds + 15% Pred 12h + 15% Proximity + 10% PMJAY.",
        "3. Live Ambulance Routing → System assigns nearest ambulance, renders GIS route, and transmits pre-arrival telemetry.",
        "4. Bed Status Toggle & Handover → Hospital staff marks bed occupied/free via 1-click non-HMS toggle (<30s).",
        "5. WebSocket Broadcast & Recalculation → Live updates propagate instantly (<100ms) to State Command Center & Patient views."
    ]
    for step in flow_steps:
        p = wf_tf.add_paragraph()
        p.text = step
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_WHITE
        p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 4: Feasibility & Viability
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Feasibility, Viability & Operational Adoption", "Overcoming real-world healthcare constraints, non-HMS facilities, and regulatory standards")

    cards_s4 = [
        ("TECHNICAL FEASIBILITY", [
            ("Low-Latency Inference: ", "Lightweight Scikit-Learn models execute predictions in <15ms per patient record."),
            ("High-Throughput WebSockets: ", "Async FastAPI backend broadcasts state deltas with sub-100ms latency."),
            ("Graceful Degradation: ", "Web Portal → Mobile PWA → 2G USSD (*999#) and SMS query fallback.")
        ], C_TEAL),
        ("OPERATIONAL ADOPTION (NON-HMS)", [
            ("Zero-Barrier Onboarding: ", "Hospitals without complex HMS utilize the Non-HMS Staff Quick Toggle Mode."),
            ("Shift Handover Batching: ", "Ward nurses update entire wards in under 30 seconds during shift changes."),
            ("IoT Auto-Sensing: ", "Smart sensors on cryogenic tanks and bed load cells eliminate manual entry errors.")
        ], C_INDIGO),
        ("REGULATORY & NATIONAL ALIGNMENT", [
            ("MeghRaj (GI Cloud) Ready: ", "Containerized microservices deployable directly on NIC Government of India Cloud."),
            ("ABDM Standards: ", "Produces standard HL7 FHIR R4 document bundles linked to 14-digit citizen ABHA IDs."),
            ("Tamper-Proof Audit: ", "SHA-256 cryptographic hash chaining prevents resource manipulation and hoarding.")
        ], C_EMERALD)
    ]

    for idx, (title, items, col) in enumerate(cards_s4):
        cx = Inches(0.8 + idx * 3.98)
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.5), Inches(3.8), Inches(5.4))
        c.fill.solid()
        c.fill.fore_color.rgb = C_CARD_BG
        c.line.color.rgb = C_CARD_BORDER
        
        ctf = c.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = Inches(0.2)
        
        p = ctf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        
        for bold_lbl, desc in items:
            p = ctf.add_paragraph()
            p.space_before = Pt(12)
            rb = p.add_run()
            rb.text = "• " + bold_lbl
            rb.font.bold = True
            rb.font.size = Pt(10.5)
            rb.font.color.rgb = C_TEXT_MAIN
            rn = p.add_run()
            rn.text = desc
            rn.font.size = Pt(10.5)
            rn.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------
    # SLIDE 5: Impact, Benefits & Rural Inclusivity
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Measurable Impact & National Health Benefits", "Quantifiable reductions in emergency transit mortality, hospital turnover gains, and rural access")

    kpi_impacts = [
        ("35–45%", "Reduction in Golden Hour Loss", "By eliminating blind hospital bouncing via specialty-aware routing", C_CORAL),
        ("20–30%", "Increase in Bed Turnover", "Enabled by 12–24h discharge foresight for elective planning", C_TEAL),
        ("< 10%", "ICU Shortage Threshold", "Automated early warning flags trigger inter-district reallocation", C_INDIGO),
        ("65%+", "Rural Citizen Inclusivity", "Full offline access via USSD (*999#) & SMS on 2G feature phones", C_EMERALD)
    ]

    for idx, (val, title, sub, col) in enumerate(kpi_impacts):
        bx = Inches(0.8 + idx * 2.98)
        b = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(1.5), Inches(2.8), Inches(1.9))
        b.fill.solid()
        b.fill.fore_color.rgb = C_CARD_BG
        b.line.color.rgb = col
        b.line.width = Pt(1.5)
        
        btf = b.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_top = Inches(0.15)
        
        p1 = btf.paragraphs[0]
        p1.text = val
        p1.font.name = "Arial"
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = btf.add_paragraph()
        p2.text = title
        p2.font.name = "Arial"
        p2.font.size = Pt(10.5)
        p2.font.bold = True
        p2.font.color.rgb = C_TEXT_MAIN
        
        p3 = btf.add_paragraph()
        p3.text = sub
        p3.font.name = "Arial"
        p3.font.size = Pt(8.5)
        p3.font.color.rgb = C_TEXT_MUTED
        p3.space_before = Pt(2)

    # Detailed Impact Columns Bottom
    bot_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.733), Inches(3.3))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = C_CARD_BG
    bot_card.line.color.rgb = C_CARD_BORDER
    
    bot_tf = bot_card.text_frame
    bot_tf.word_wrap = True
    bot_tf.margin_left = bot_tf.margin_top = Inches(0.2)
    
    bp = bot_tf.paragraphs[0]
    bp.text = "SOCIETAL, CLINICAL & ADMINISTRATIVE BENEFITS"
    bp.font.name = "Arial"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = C_TEAL

    benefit_points = [
        ("Clinical Triage & Golden Hour: ", "Specialty matching ensures stroke, cardiac, and trauma emergencies arrive at equipped facilities on the first attempt."),
        ("Hospital Stress Relief: ", "Replaces chaotic telephone inquiries with automated digital capacity broadcasting, reducing clinician administrative burden."),
        ("Epidemic & Disaster Preparedness: ", "The Hospital Digital Twin simulates surge scenarios (+20%, +50%), calculating exact hours to ICU depletion and recommending step-down protocols."),
        ("Rural India Inclusion: ", "Ensures non-smartphone citizens dialling *999# or sending SMS queries receive instant localized emergency hospital recommendations.")
    ]
    for bold_lbl, desc in benefit_points:
        p = bot_tf.add_paragraph()
        p.space_before = Pt(6)
        rb = p.add_run()
        rb.text = "✔ " + bold_lbl
        rb.font.bold = True
        rb.font.size = Pt(10.5)
        rb.font.color.rgb = C_TEXT_MAIN
        rn = p.add_run()
        rn.text = desc
        rn.font.size = Pt(10.5)
        rn.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------
    # SLIDE 6: Roadmap & Implementation
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Research, Dependencies & 3-Phase Roadmap", "Clinical feature selection, external API integrations, and structured rollout strategy")

    # Table for 3-Phase Roadmap
    rows, cols = 4, 3
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.0))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(7.933)

    headers = ["ROADMAP PHASE", "TIMELINE", "KEY DELIVERABLES & MILESTONES"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

    roadmap_data = [
        ("Phase 1: Working Prototype\n(COMPLETED)", "Month 1–2", "• Full-stack FastAPI + React prototype with 100% passing test suite.\n• Pre-trained Scikit-Learn Predictive Bed Turnover ML engine.\n• 16 hospitals seed dataset, real-time WebSockets, Leaflet GPS & USSD simulator."),
        ("Phase 2: District Pilot\n(Pilot Validation)", "Month 3–4", "• Pilot rollout across 1 district (e.g. Vellore / Chennai) connecting 5 GHs & private HMS.\n• Integration with state 108 Emergency Ambulance dispatch systems.\n• Real-world calibration of Length-of-Stay prediction models."),
        ("Phase 3: Statewide Scaling\n(ABDM Production)", "Month 5–6", "• Full deployment on MeghRaj (GI Cloud) integrated with State Health Dept.\n• Official ABDM M1/M2/M3 milestone certification & ABHA linking.\n• National rollout of BSNL/DoT rural USSD (*999#) gateway.")
    ]

    for row_idx, data in enumerate(roadmap_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if row_idx % 2 == 1 else C_WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.name = "Arial"
            p.font.size = Pt(9.5)
            p.font.color.rgb = C_TEXT_MAIN

    # Dependencies & Risks Bottom Card
    dep_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.733), Inches(2.2))
    dep_card.fill.solid()
    dep_card.fill.fore_color.rgb = C_CARD_BG
    dep_card.line.color.rgb = C_CARD_BORDER
    dep_tf = dep_card.text_frame
    dep_tf.word_wrap = True
    dep_tf.margin_left = dep_tf.margin_top = Inches(0.15)
    
    p = dep_tf.paragraphs[0]
    p.text = "RESEARCH BACKING & DEPENDENCIES"
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_INDIGO

    deps = [
        ("Clinical Feature Backing: ", "Evaluates 10 features: age, diagnosis category, stay hours, clinical stage (ICU Critical/Step-down/Oral meds), SpO2, HR, MAP, RR, Temp, stability index."),
        ("External Gateway Integrations: ", "ABDM Sandbox APIs for ABHA linking; National DoT/BSNL telecom gateway for USSD (*999#); OpenStreetMap / MapmyIndia for localized routing."),
        ("Risk Mitigation: ", "Non-HMS toggle mode solves low HMS adoption; offline SMS/USSD resolves rural connectivity dropouts; SHA-256 audit log eliminates data tampering.")
    ]
    for b_lbl, desc in deps:
        p = dep_tf.add_paragraph()
        p.space_before = Pt(4)
        rb = p.add_run()
        rb.text = "• " + b_lbl
        rb.font.bold = True
        rb.font.size = Pt(10)
        rb.font.color.rgb = C_TEXT_MAIN
        rn = p.add_run()
        rn.text = desc
        rn.font.size = Pt(10)
        rn.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------
    # SLIDE 7: Conclusion & Summary
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = C_DARK
    bg7.line.color.rgb = C_DARK

    title7 = slide7.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11.5), Inches(1.2))
    tf7 = title7.text_frame
    p = tf7.paragraphs[0]
    p.text = "Why MedFlow Wins — Summary"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    p2 = tf7.add_paragraph()
    p2.text = "Converting reactive hospital resource management into proactive, predictive national coordination"
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = C_TEAL
    p2.space_before = Pt(4)

    # 4 Key Takeaway Cards
    takeaways = [
        ("1. SYSTEM-LEVEL THINKING", "Addresses the entire healthcare chain: Patients (search/referral), Hospitals (toggles/ML turnover), and Government (statewide heatmap & reallocation).", C_TEAL),
        ("2. CLEAR ML DIFFERENTIATOR", "12–24h Predictive Bed Turnover Engine with clinical explainability weights outperforms static 'bed counter' apps.", C_INDIGO),
        ("3. INDIA-FIRST DESIGN", "Integrated with ABDM/FHIR R4, MeghRaj cloud readiness, Ayushman Bharat (PMJAY), and 2G USSD (*999#) rural fallback.", C_EMERALD),
        ("4. WORKING MVP READY", "Fully functional prototype with FastAPI, React, Scikit-Learn ML, WebSockets, and 12/12 passing test suite.", C_CORAL)
    ]

    for idx, (head, body, col) in enumerate(takeaways):
        bx = Inches(0.9 + idx * 2.9)
        b = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(2.3), Inches(2.7), Inches(3.6))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(30, 41, 59)
        b.line.color.rgb = col
        b.line.width = Pt(1.5)
        
        btf = b.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_top = Inches(0.18)
        
        bp1 = btf.paragraphs[0]
        bp1.text = head
        bp1.font.name = "Arial"
        bp1.font.size = Pt(11)
        bp1.font.bold = True
        bp1.font.color.rgb = col
        
        bp2 = btf.add_paragraph()
        bp2.text = body
        bp2.font.name = "Arial"
        bp2.font.size = Pt(10.5)
        bp2.font.color.rgb = C_WHITE
        bp2.space_before = Pt(8)

    # Bottom Banner
    banner = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.7))
    banner.fill.solid()
    banner.fill.fore_color.rgb = C_TEAL_DARK
    banner.line.color.rgb = C_TEAL
    banner_tf = banner.text_frame
    banner_p = banner_tf.paragraphs[0]
    banner_p.text = "MedFlow: Empowering India's Healthcare Ecosystem with Real-Time Predictive Intelligence"
    banner_p.alignment = PP_ALIGN.CENTER
    banner_p.font.name = "Arial"
    banner_p.font.size = Pt(12)
    banner_p.font.bold = True
    banner_p.font.color.rgb = C_TEAL

    output_path = "MedFlow_VITISH_2026_Idea_Submission.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
